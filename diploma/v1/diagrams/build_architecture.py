# -*- coding: utf-8 -*-
"""Native PowerPoint Architecture diagram on slide 8 (classic, monochrome, full-page).

Показывает реальную архитектуру мигрированного проекта:
конвейер подготовки контента → бессерверный backend (Firebase) →
клиент React + TypeScript + Vite, организованный по методологии Feature-Sliced Design.
"""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BASE = r"C:\My Projects\E-Magios-Core-Site\diploma\v1"
PPTX = os.environ.get("PPTX_OUT", os.path.join(BASE, "221_3711_ГайдарьМД_Диплом_Презентация_v3.pptx"))

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT  = "Arial"

def E(inch): return Emu(int(inch * 914400))

prs    = Presentation(PPTX)
slide  = prs.slides[7]   # slide 8 (0-indexed)
shapes = slide.shapes

# ── CLEANUP: remove old architecture PNG (image16) ──────────────
for sh in list(shapes):
    if sh.shape_type == 13:
        try:
            rId  = sh._element.blip_rId
            part = slide.part.related_part(rId)
            if "image16" in part.partname:
                sh._element.getparent().remove(sh._element)
        except Exception:
            pass

# ── HELPERS ──────────────────────────────────────────────────────

def _line(shape, width=1.0, dash=False):
    shape.line.color.rgb = BLACK
    shape.line.width = Pt(width)
    ln = shape.line._get_or_add_ln()
    e = ln.find(qn("a:prstDash"))
    if e is not None:
        ln.remove(e)
    if dash:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    return ln


def box(cx, cy, w, h, text, fs=11, bold=False, multiline=True):
    s = shapes.add_shape(MSO_SHAPE.RECTANGLE, E(cx - w/2), E(cy - h/2), E(w), E(h))
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    _line(s, 1.25)
    s.shadow.inherit = False
    tf = s.text_frame
    tf.word_wrap = multiline
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = E(0.06)
    tf.margin_top  = tf.margin_bottom = E(0.03)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.name = FONT
    r.font.bold = bold; r.font.color.rgb = BLACK
    return {"cx": cx, "cy": cy, "w": w, "h": h}


def band(cx, cy, w, h, name, desc, fs=10):
    """FSD layer band: bold layer name + plain contents, left-aligned."""
    s = shapes.add_shape(MSO_SHAPE.RECTANGLE, E(cx - w/2), E(cy - h/2), E(w), E(h))
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    _line(s, 1.10)
    s.shadow.inherit = False
    tf = s.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = E(0.12); tf.margin_right = E(0.06)
    tf.margin_top  = tf.margin_bottom = E(0.02)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run(); r1.text = name + "  —  "
    r1.font.size = Pt(fs); r1.font.name = FONT; r1.font.bold = True; r1.font.color.rgb = BLACK
    r2 = p.add_run(); r2.text = desc
    r2.font.size = Pt(fs); r2.font.name = FONT; r2.font.color.rgb = BLACK
    return {"cx": cx, "cy": cy, "w": w, "h": h}


def cluster(x, y, w, h, label, fs=11):
    """Cluster border (sent to back) with label at top-left."""
    s = shapes.add_shape(MSO_SHAPE.RECTANGLE, E(x), E(y), E(w), E(h))
    s.fill.background()
    _line(s, 1.10)
    s.shadow.inherit = False
    tf = s.text_frame; tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_top  = E(0.05)
    tf.margin_left = E(0.12)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = label
    r.font.size = Pt(fs); r.font.name = FONT
    r.font.bold = True; r.font.color.rgb = BLACK
    sp = s._element; sp.getparent().remove(sp)
    shapes._spTree.insert(2, sp)
    return s


def arrow(x1, y1, x2, y2, label="", fs=9, dash=False, lx_off=0, ly_off=-0.14):
    c = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, E(x1), E(y1), E(x2), E(y2))
    ln = _line(c, 1.25, dash)
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    c.shadow.inherit = False
    if label:
        lx = (x1 + x2) / 2 + lx_off
        ly = (y1 + y2) / 2 + ly_off
        tw, th = 1.80, 0.22
        tb = shapes.add_textbox(E(lx - tw/2), E(ly - th/2), E(tw), E(th))
        tf = tb.text_frame; tf.word_wrap = False
        for m in ("margin_left","margin_right","margin_top","margin_bottom"):
            setattr(tf, m, 0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label
        r.font.size = Pt(fs); r.font.name = FONT
        r.font.italic = True; r.font.color.rgb = BLACK


def edges(b, side):
    """Shorthand for box edge coordinates."""
    cx, cy, w, h = b["cx"], b["cy"], b["w"], b["h"]
    if side == "l": return cx - w/2, cy
    if side == "r": return cx + w/2, cy
    if side == "t": return cx, cy - h/2
    if side == "b": return cx, cy + h/2
    if side == "tl": return cx - w/2, cy - h/2
    if side == "tr": return cx + w/2, cy - h/2
    if side == "bl": return cx - w/2, cy + h/2
    if side == "br": return cx + w/2, cy + h/2


# ════════════════════════════════════════════════════════════════
# LAYOUT  (slide = 10" × 7.5"; контент y ≈ 1.50"–6.80")
# ════════════════════════════════════════════════════════════════
LM, RM = 0.30, 9.70
FULLW  = RM - LM

b = {}

# ── 1. Конвейер подготовки контента (офлайн) ─────────────────────
cluster(LM, 1.50, FULLW, 0.98, "Подготовка контента (офлайн): Obsidian → JSON → Firestore")
PY = 2.08
b["ov"]   = box(1.50, PY, 2.00, 0.60, "Obsidian Vault\n(Markdown)", fs=10)
b["pars"] = box(3.80, PY, 2.10, 0.60, "Python-парсеры\n(14 шт., MD → JSON)", fs=10)
b["val"]  = box(6.05, PY, 1.95, 0.60, "Валидация и\nнормализация", fs=10)
b["imp"]  = box(8.35, PY, 2.05, 0.60, "Import CLI\n(Firebase Admin SDK)", fs=10)
arrow(*edges(b["ov"],  "r"), *edges(b["pars"], "l"))
arrow(*edges(b["pars"],"r"), *edges(b["val"],  "l"))
arrow(*edges(b["val"], "r"), *edges(b["imp"],  "l"))

# ── 2. Firebase (бессерверный backend) ───────────────────────────
cluster(LM, 2.66, FULLW, 1.30, "Firebase — бессерверный backend")
FY = 3.42
b["auth"]  = box(1.85, FY, 2.55, 0.82,
                 "Firebase Authentication\n(Google OAuth 2.0,\ncustom claim: role)", fs=9.5)
b["fs"]    = box(5.00, FY, 3.00, 0.82,
                 "Cloud Firestore\n14 коллекций справочника · news ·\nusers/characters/diceRolls · contentManifest", fs=8.5)
b["rules"] = box(8.20, FY, 2.55, 0.82,
                 "Security Rules\n(доступ по ролям author/\neditor/admin и по статусам)", fs=9.5)
# импорт из конвейера в Firestore
arrow(*edges(b["imp"], "b"), *edges(b["fs"], "t"), label="импорт", lx_off=0.55, ly_off=-0.12)
# Security Rules ограничивают доступ к Firestore (подпись не нужна — смысл из текста блока)
arrow(*edges(b["rules"], "l"), *edges(b["fs"], "r"), dash=True)

# ── 3. Клиент: React + TypeScript + Vite (Feature-Sliced Design) ─
cluster(LM, 4.46, FULLW, 2.32,
        "Клиент · React + TypeScript + Vite · Feature-Sliced Design   (зависимости — сверху вниз)")
BX, BW, BH = 5.00, 9.00, 0.305
y0 = 4.84
gap = 0.025
layers = [
    ("app",      "маршрутизация (React Router) · провайдеры (Auth, Dice) · телеметрия переходов"),
    ("pages",    "Главная · Новости · База данных · Редактор персонажей · Dashboard · Профиль · Книга"),
    ("widgets",  "Layout · Sidebar · Таблица справочника · Панель фильтров · Модальные карточки"),
    ("features", "авторизация · фильтры и сортировка БД · редактор персонажей · броски кубов · аналитика"),
    ("entities", "compendium · character · content · user   (типы · DTO · mappers · schema-валидация)"),
    ("shared",   "Repository Layer · кэш IndexedDB · Firebase-клиент · UI-кит · телеметрия"),
]
rows = {}
for i, (name, desc) in enumerate(layers):
    cy = y0 + i * (BH + gap)
    rows[name] = band(BX, cy, BW, BH, name, desc, fs=9.5)

# ── 4. Потоки данных между Firebase и клиентом (короткие, в зазоре) ─
# Cloud Firestore → клиент: чтение опубликованного контента + запись по ролям
arrow(5.80, edges(b["fs"], "b")[1], 6.45, 4.46,
      label="чтение published · запись по ролям", fs=8.5, lx_off=1.05, ly_off=-0.04)
# Firebase Auth → клиент: сессия и роль (custom claim) — вход в левую грань кластера
arrow(edges(b["auth"], "b")[0], edges(b["auth"], "b")[1], LM, 5.05,
      label="сессия · роль", fs=8.5, dash=True, lx_off=-0.55, ly_off=-0.10)

# ──────────────────────────────────────────────────────────────────
prs.save(PPTX)
print("ARCH BUILT ->", PPTX)
