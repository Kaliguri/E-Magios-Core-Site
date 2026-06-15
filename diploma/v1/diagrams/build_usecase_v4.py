# -*- coding: utf-8 -*-
"""Rebuild the Use Case on slide 6 of v4 — simplified per supervisor feedback.

Changes vs. the old version:
  * Cloud Firestore removed from actors (it is infrastructure, not an actor —
    it belongs in the architecture/DFD slide). External actors left: Google
    Authentication, Obsidian Vault.
  * Roles condensed: Гость -> чтение/поиск; Пользователь -> профиль/персонажи/
    кубы (3 самостоятельных варианта, без вложенных include); Админ ->
    импорт/публикация контента.
  * No «include» chains, no info panel -> диаграмма не перегружена.

Edits v4 DIRECTLY (not the integrate->pack pipeline): first removes the old
diagram shapes (non-placeholder, top >= 1.40), keeping title/subtitle/page no.,
then redraws.
"""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BASE = r"C:\My Projects\E-Magios-Core-Site\diploma\v1"
PPTX = os.environ.get("PPTX_OUT",
                      os.path.join(BASE, "221_3711_ГайдарьМД_Диплом_Презентация_v4.pptx"))

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT  = "Arial"

def E(inch):
    return Emu(int(inch * 914400))

prs    = Presentation(PPTX)
slide  = prs.slides[5]            # slide 6 (0-indexed)
shapes = slide.shapes

# ── CLEANUP: remove old diagram shapes, keep header/placeholders ──
for sh in list(shapes):
    if sh.is_placeholder:
        continue
    top = sh.top / 914400 if sh.top is not None else -99
    if top >= 1.40:
        sh._element.getparent().remove(sh._element)

# ── HELPERS ──────────────────────────────────────────────────────

def _line(shape, width=1.0, dash=False):
    shape.line.color.rgb = BLACK
    shape.line.width = Pt(width)
    ln = shape.line._get_or_add_ln()
    e = ln.find(qn("a:prstDash"))
    if e is not None:
        ln.remove(e)
    if dash:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    return ln


def oval(cx, cy, w, h, text, fs=12):
    s = shapes.add_shape(MSO_SHAPE.OVAL, E(cx - w/2), E(cy - h/2), E(w), E(h))
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    _line(s, 1.0)
    s.shadow.inherit = False
    tf = s.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ("margin_left", "margin_right"):
        setattr(tf, m, E(0.04))
    for m in ("margin_top", "margin_bottom"):
        setattr(tf, m, E(0.01))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.name = FONT; r.font.color.rgb = BLACK
    return {"cx": cx, "cy": cy, "w": w, "h": h}


def boundary_rect(x, y, w, h, title):
    s = shapes.add_shape(MSO_SHAPE.RECTANGLE, E(x), E(y), E(w), E(h))
    s.fill.background()
    _line(s, 1.50)
    s.shadow.inherit = False
    tf = s.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_top = E(0.05)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title
    r.font.size = Pt(13); r.font.name = FONT; r.font.bold = True; r.font.color.rgb = BLACK
    sp = s._element
    sp.getparent().remove(sp)
    shapes._spTree.insert(2, sp)
    return s


def connector(x1, y1, x2, y2, width=1.0, dash=False, arrow="triangle"):
    c = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, E(x1), E(y1), E(x2), E(y2))
    ln = _line(c, width, dash)
    if arrow:
        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": arrow, "w": "med", "len": "med"}))
    c.shadow.inherit = False
    return c


def label(cx, cy, text, fs=8):
    w, h = 1.2, 0.22
    tb = shapes.add_textbox(E(cx - w/2), E(cy - h/2), E(w), E(h))
    tf = tb.text_frame; tf.word_wrap = False
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.name = FONT
    r.font.italic = True; r.font.color.rgb = BLACK


def actor(cx, cy, text, wl=1.50):
    top = cy - 0.40
    hd  = 0.24
    h = shapes.add_shape(MSO_SHAPE.OVAL, E(cx - hd/2), E(top), E(hd), E(hd))
    h.fill.solid(); h.fill.fore_color.rgb = WHITE
    _line(h, 1.0); h.shadow.inherit = False
    bt = top + hd
    connector(cx, bt,            cx, bt + 0.36, arrow=None)
    connector(cx - 0.24, bt + 0.10, cx + 0.24, bt + 0.10, arrow=None)
    connector(cx, bt + 0.36, cx - 0.20, bt + 0.66, arrow=None)
    connector(cx, bt + 0.36, cx + 0.20, bt + 0.66, arrow=None)
    tb = shapes.add_textbox(E(cx - wl/2), E(top + 0.88), E(wl), E(0.62))
    tf = tb.text_frame; tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(10); r.font.name = FONT; r.font.color.rgb = BLACK
    return {"cx": cx, "cy": cy}


# ── LAYOUT (slide 10.00" × 7.50") ────────────────────────────────
# No external actors: Use Case = роли + цели; внешние системы (Google
# Auth, Obsidian) живут в DFD. Composition centred to fill the page.
CX = 5.35          # use-case column centre
OW, OH = 3.40, 0.50

# grouped rows (gaps between actor groups)
g = [2.05, 2.55, 3.05, 3.55]        # Гость (4)
u = [4.05, 4.53, 5.01]              # Пользователь (3)
a = [5.55, 6.05]                    # Администратор (2)

boundary_rect(3.30, 1.50, 4.30, 5.28, "E'Magios Core Site")

uc = {}
uc["books"]  = oval(CX, g[0], OW, OH, "Читать книги")
uc["news"]   = oval(CX, g[1], OW, OH, "Смотреть новости")
uc["search"] = oval(CX, g[2], OW, OH, "Искать по базе данных")
uc["login"]  = oval(CX, g[3], OW, OH, "Войти через Google")

uc["profile"] = oval(CX, u[0], OW, OH, "Открыть профиль")
uc["chars"]   = oval(CX, u[1], OW, OH, "Работать с персонажами")
uc["dice"]    = oval(CX, u[2], OW, OH, "Бросать кубы")

uc["import"]  = oval(CX, a[0], OW, OH, "Импортировать контент")
uc["publish"] = oval(CX, a[1], OW, OH, "Публиковать контент")

# actors (left) — only the three human roles
gost  = actor(1.45, (g[0] + g[3]) / 2, "Гость", wl=1.20)
user  = actor(1.45, (u[0] + u[2]) / 2, "Пользователь", wl=1.40)
admin = actor(1.45, (a[0] + a[1]) / 2, "Администратор контента", wl=1.60)

# ── CONNECTORS ───────────────────────────────────────────────────

def lft(k): return (uc[k]["cx"] - uc[k]["w"] / 2, uc[k]["cy"])
def rgt(k): return (uc[k]["cx"] + uc[k]["w"] / 2, uc[k]["cy"])


def assoc(act, k):
    x1, y1 = act["cx"] + 0.30, act["cy"]
    x2, y2 = lft(k)
    connector(x1, y1, x2 - 0.02, y2, arrow=None)


for k in ("books", "news", "search", "login"):
    assoc(gost, k)
for k in ("profile", "chars", "dice"):
    assoc(user, k)
for k in ("import", "publish"):
    assoc(admin, k)

prs.save(PPTX)
print("USE CASE v4 BUILT ->", PPTX)
