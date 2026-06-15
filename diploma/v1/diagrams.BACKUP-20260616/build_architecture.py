# -*- coding: utf-8 -*-
"""Native PowerPoint Architecture diagram on slide 8 (classic, monochrome, full-page)."""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BASE = r"C:\My Projects\E-Magios-Core-Site\diploma\v1"
PPTX = os.path.join(BASE, "221_3711_ГайдарьМД_Диплом_Презентация_v3.pptx")

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT  = "Arial"

def E(inch): return Emu(int(inch * 914400))

prs    = Presentation(PPTX)
slide  = prs.slides[7]   # slide 8 (0-indexed)
shapes = slide.shapes

# ── CLEANUP: remove old architecture PNG (image16) ──────────────
for sh in list(shapes):
    if sh.shape_type == 13:
        try:
            rId  = sh._element.blip_rId
            part = slide.part.related_part(rId)
            if "image16" in part.partname:
                sh._element.getparent().remove(sh._element)
        except Exception:
            pass

# ── HELPERS ──────────────────────────────────────────────────────

def _line(shape, width=1.0, dash=False):
    shape.line.color.rgb = BLACK
    shape.line.width = Pt(width)
    ln = shape.line._get_or_add_ln()
    e = ln.find(qn("a:prstDash"))
    if e is not None:
        ln.remove(e)
    if dash:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    return ln


def box(cx, cy, w, h, text, fs=11, bold=False, multiline=True):
    s = shapes.add_shape(MSO_SHAPE.RECTANGLE, E(cx - w/2), E(cy - h/2), E(w), E(h))
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    _line(s, 1.25)
    s.shadow.inherit = False
    tf = s.text_frame
    tf.word_wrap = multiline
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = E(0.06)
    tf.margin_top  = tf.margin_bottom = E(0.04)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.name = FONT
    r.font.bold = bold; r.font.color.rgb = BLACK
    return {"cx": cx, "cy": cy, "w": w, "h": h}


def cluster(x, y, w, h, label, fs=11):
    """Cluster border (sent to back) with label at top-left."""
    s = shapes.add_shape(MSO_SHAPE.RECTANGLE, E(x), E(y), E(w), E(h))
    s.fill.background()
    _line(s, 1.10)
    s.shadow.inherit = False
    tf = s.text_frame; tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_top  = E(0.05)
    tf.margin_left = E(0.12)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = label
    r.font.size = Pt(fs); r.font.name = FONT
    r.font.bold = True; r.font.color.rgb = BLACK
    sp = s._element; sp.getparent().remove(sp)
    shapes._spTree.insert(2, sp)
    return s


def arrow(x1, y1, x2, y2, label="", fs=9, dash=False, lx_off=0, ly_off=-0.14):
    c = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, E(x1), E(y1), E(x2), E(y2))
    ln = _line(c, 1.25, dash)
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    c.shadow.inherit = False
    if label:
        lx = (x1 + x2) / 2 + lx_off
        ly = (y1 + y2) / 2 + ly_off
        tw, th = 1.50, 0.22
        tb = shapes.add_textbox(E(lx - tw/2), E(ly - th/2), E(tw), E(th))
        tf = tb.text_frame; tf.word_wrap = False
        for m in ("margin_left","margin_right","margin_top","margin_bottom"):
            setattr(tf, m, 0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label
        r.font.size = Pt(fs); r.font.name = FONT
        r.font.italic = True; r.font.color.rgb = BLACK


def edges(b, side):
    """Shorthand for box edge coordinates."""
    cx, cy, w, h = b["cx"], b["cy"], b["w"], b["h"]
    if side == "l": return cx - w/2, cy
    if side == "r": return cx + w/2, cy
    if side == "t": return cx, cy - h/2
    if side == "b": return cx, cy + h/2
    if side == "tl": return cx - w/2, cy - h/2
    if side == "tr": return cx + w/2, cy - h/2
    if side == "bl": return cx - w/2, cy + h/2
    if side == "br": return cx + w/2, cy + h/2


# ── LAYOUT (slide = 10" × 7.5", content y=1.55"–6.80") ──────────
BW = 2.15   # pipeline box width
BH = 0.78   # pipeline box height
PY = 2.19   # pipeline row y-center

# ── 1. Content pipeline cluster ──────────────────────────────────
cluster(0.28, 1.55, 9.44, 1.26, "Подготовка контента (офлайн)")

b = {}
b["ov"]   = box(1.43,  PY, BW,   BH, "Obsidian Vault\n(Markdown)")
b["pars"] = box(3.73,  PY, BW,   BH, "Python-парсеры\n(Markdown → JSON)")
b["json"] = box(5.97,  PY, 1.90, BH, "JSON-файлы")
b["imp"]  = box(8.22,  PY, BW,   BH, "Import-скрипт\n(Firebase Admin SDK)")

# arrows inside pipeline (horizontal, L→R)
arrow(*edges(b["ov"],   "r"), *edges(b["pars"], "l"))
arrow(*edges(b["pars"], "r"), *edges(b["json"], "l"))
arrow(*edges(b["json"], "r"), *edges(b["imp"],  "l"))

# ── 2. Cloud Firestore (wide box, center) ────────────────────────
b["fs"] = box(5.0, 3.90, 5.20, 1.05,
              "Cloud Firestore\n(справочник, новости, профили, персонажи, манифест версий)",
              fs=10)

# Import → Firestore (diagonal down-left)
arrow(*edges(b["imp"], "b"),
      edges(b["fs"], "r")[0], edges(b["fs"], "t")[1],
      label="импорт", lx_off=0.35, ly_off=-0.14)

# ── 3. Client app cluster ─────────────────────────────────────────
cluster(0.28, 4.86, 6.70, 1.88, "Клиент: React + TypeScript + Vite")

b["repo"] = box(2.10, 5.50, 2.60, 0.75,
                "Repository Layer\n(доступ к данным,\nфильтр по статусу)", fs=10)
b["idb"]  = box(1.35, 6.43, 2.10, 0.75, "IndexedDB\n(кэш справочника)", fs=10)
b["ui"]   = box(4.85, 6.43, 2.50, 0.75, "React UI\n(книги, база, редактор)", fs=10)

# ── 4. External: Firebase Auth + User ────────────────────────────
b["auth"] = box(8.65, 5.50, 2.20, 0.75, "Firebase\nAuthentication")
b["user"] = box(8.65, 6.43, 2.20, 0.75, "Пользователь")

# ── CONNECTORS ───────────────────────────────────────────────────

# Firestore → Repository (diagonal down-left)
arrow(edges(b["fs"], "l")[0] + 0.60, edges(b["fs"], "b")[1],
      *edges(b["repo"], "t"),
      label="чтение", lx_off=-0.45, ly_off=-0.14)

# Repository → IndexedDB
arrow(b["repo"]["cx"] - 0.35, edges(b["repo"], "b")[1],
      *edges(b["idb"], "t"),
      label="кэш", lx_off=-0.38, ly_off=-0.12)

# Repository → React UI
arrow(b["repo"]["cx"] + 0.35, edges(b["repo"], "b")[1],
      *edges(b["ui"], "t"),
      label="данные", lx_off=0.30, ly_off=-0.12)

# Firebase Auth → React UI (dashed)
arrow(*edges(b["auth"], "l"),
      *edges(b["ui"],   "r"),
      label="вход через Google", fs=9, dash=True, ly_off=-0.16)

# User → React UI (horizontal)
arrow(*edges(b["user"], "l"),
      *edges(b["ui"],   "r"),
      label="работа с сайтом", fs=9, ly_off=+0.14)

# ──────────────────────────────────────────────────────────────────
prs.save(PPTX)
print("ARCH BUILT ->", PPTX)
