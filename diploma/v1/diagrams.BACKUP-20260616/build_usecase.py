# -*- coding: utf-8 -*-
"""Native PowerPoint Use Case diagram on slide 6 (classic, monochrome, full-width + info panel)."""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BASE = r"C:\My Projects\E-Magios-Core-Site\diploma\v1"
PPTX = os.path.join(BASE, "221_3711_ГайдарьМД_Диплом_Презентация_v3.pptx")

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT  = "Arial"

def E(inch):
    return Emu(int(inch * 914400))

prs   = Presentation(PPTX)
slide = prs.slides[5]          # slide 6 (0-indexed)
shapes = slide.shapes

# ================================================================
# CLEANUP
# ================================================================

# 1. Fix subtitle typo ("Фнкциональные" -> "Функциональные")
for sh in list(shapes):
    if sh.has_text_frame:
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if "Фнкциональные" in run.text:
                    run.text = run.text.replace("Фнкциональные", "Функциональные")

# 2. Remove old use-case picture (image15 from v2)
for sh in list(shapes):
    if sh.shape_type == 13:   # PICTURE
        try:
            rId  = sh._element.blip_rId
            part = slide.part.related_part(rId)
            if "image15" in part.partname:
                sh._element.getparent().remove(sh._element)
        except Exception:
            pass

# 3. Remove legend shapes ("Условные обозначения" block from v2)
LEGEND_KEYWORDS = ("Условные", "Сплошная", "Пунктирная", "Стрелка-треугольник")
for sh in list(shapes):
    if sh.has_text_frame:
        txt = sh.text_frame.text
        if any(kw in txt for kw in LEGEND_KEYWORDS):
            sh._element.getparent().remove(sh._element)

# ================================================================
# HELPERS
# ================================================================

def _line(shape, width=1.0, dash=False):
    shape.line.color.rgb = BLACK
    shape.line.width = Pt(width)
    ln = shape.line._get_or_add_ln()
    e = ln.find(qn("a:prstDash"))
    if e is not None:
        ln.remove(e)
    if dash:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    return ln


def oval(cx, cy, w, h, text, fs=10):
    s = shapes.add_shape(MSO_SHAPE.OVAL, E(cx - w/2), E(cy - h/2), E(w), E(h))
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    _line(s, 1.0)
    s.shadow.inherit = False
    tf = s.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ("margin_left", "margin_right"):
        setattr(tf, m, E(0.04))
    for m in ("margin_top", "margin_bottom"):
        setattr(tf, m, E(0.01))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.name = FONT; r.font.color.rgb = BLACK
    return {"cx": cx, "cy": cy, "w": w, "h": h}


def boundary_rect(x, y, w, h, title):
    """System boundary rectangle — sent to back so ovals draw on top."""
    s = shapes.add_shape(MSO_SHAPE.RECTANGLE, E(x), E(y), E(w), E(h))
    s.fill.background()
    _line(s, 1.50)
    s.shadow.inherit = False
    tf = s.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_top = E(0.04)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title
    r.font.size = Pt(13); r.font.name = FONT; r.font.bold = True; r.font.color.rgb = BLACK
    sp = s._element
    sp.getparent().remove(sp)
    shapes._spTree.insert(2, sp)
    return s


def connector(x1, y1, x2, y2, width=1.0, dash=False, arrow="triangle"):
    c = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, E(x1), E(y1), E(x2), E(y2))
    ln = _line(c, width, dash)
    if arrow:
        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": arrow, "w": "med", "len": "med"}))
    c.shadow.inherit = False
    return c


def label(cx, cy, text, fs=8):
    w, h = 1.0, 0.22
    tb = shapes.add_textbox(E(cx - w/2), E(cy - h/2), E(w), E(h))
    tf = tb.text_frame; tf.word_wrap = False
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.name = FONT
    r.font.italic = True; r.font.color.rgb = BLACK


def actor(cx, cy, text, wl=1.40):
    """Draw UML actor stickman + label."""
    top = cy - 0.40
    hd  = 0.22
    h = shapes.add_shape(MSO_SHAPE.OVAL, E(cx - hd/2), E(top), E(hd), E(hd))
    h.fill.solid(); h.fill.fore_color.rgb = WHITE
    _line(h, 1.0); h.shadow.inherit = False
    bt = top + hd
    connector(cx, bt,       cx, bt + 0.34, arrow=None)           # spine
    connector(cx - 0.22, bt + 0.09, cx + 0.22, bt + 0.09, arrow=None)  # arms
    connector(cx, bt + 0.34, cx - 0.18, bt + 0.62, arrow=None)  # leg L
    connector(cx, bt + 0.34, cx + 0.18, bt + 0.62, arrow=None)  # leg R
    tb = shapes.add_textbox(E(cx - wl/2), E(top + 0.82), E(wl), E(0.65))
    tf = tb.text_frame; tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(9); r.font.name = FONT; r.font.color.rgb = BLACK
    return {"cx": cx, "cy": cy}


def info_panel(x, y, w, h):
    """Right-side info box — replaces the legend, summarises the project."""
    s = shapes.add_shape(MSO_SHAPE.RECTANGLE, E(x), E(y), E(w), E(h))
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    _line(s, 1.25)
    s.shadow.inherit = False

    tf = s.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left   = E(0.10)
    tf.margin_right  = E(0.08)
    tf.margin_top    = E(0.09)
    tf.margin_bottom = E(0.05)

    def row(text, bold=False, fs=8.5, align=PP_ALIGN.LEFT):
        p = tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size      = Pt(fs)
        r.font.bold      = bold
        r.font.name      = FONT
        r.font.color.rgb = BLACK

    # --- header ---
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run(); r0.text = "О ПРОЕКТЕ"
    r0.font.size = Pt(11); r0.font.bold = True
    r0.font.name = FONT;   r0.font.color.rgb = BLACK

    row("")
    row("Роли пользователей:",  bold=True, fs=9)
    row("Гость — просмотр книг,")
    row("новостей, базы данных")
    row("")
    row("Авторизованный —")
    row("профиль + редактор")
    row("персонажей")
    row("")
    row("Администратор —")
    row("импорт данных из")
    row("Obsidian в Firestore")
    row("")
    row("Редактор персонажей:", bold=True, fs=9)
    row("7 CRUD-операций:")
    row("создать, редактировать,")
    row("сохранить, загрузить,")
    row("удалить, экспортировать,")
    row("импортировать")
    row("")
    row("Внешние сервисы:", bold=True, fs=9)
    row("  Google OAuth 2.0")
    row("  Cloud Firestore")
    row("  Obsidian Vault")

    return s


# ================================================================
# LAYOUT  (slide is 10.00" × 7.50" — classic 4:3)
# ================================================================

# 9 rows; first row starts at 2.00" to leave room for boundary title
r = [2.00 + i * 0.53 for i in range(9)]
# → 2.00, 2.53, 3.06, 3.59, 4.12, 4.65, 5.18, 5.71, 6.24

LX  = 2.95   # left-column oval centres (main use cases)
RX  = 6.35   # right-column oval centres (include sub-cases)
LW  = 2.50   # left oval width
RW  = 2.25   # right oval width
OH  = 0.44   # oval height

# System boundary — drawn first → sent to back
boundary_rect(0.88, 1.46, 7.72, 5.25, "E'Magios Core Site")

# ---- LEFT COLUMN: primary use cases ----
uc = {}
uc["enter"]   = oval(LX, r[0], LW, OH, "Зайти на сайт",       fs=11)
uc["books"]   = oval(LX, r[1], LW, OH, "Читать книги",         fs=11)
uc["news"]    = oval(LX, r[2], LW, OH, "Смотреть новости",     fs=11)
uc["db"]      = oval(LX, r[3], LW, OH, "Пользоваться БД",      fs=11)
uc["login"]   = oval(LX, r[4], LW, OH, "Войти через Google",   fs=11)
uc["profile"] = oval(LX, r[5], LW, OH, "Открыть профиль",      fs=11)
uc["editor"]  = oval(LX, r[6], LW, OH, "Работать с редактором персонажей", fs=9)
uc["import"]  = oval(LX, r[7], LW, OH, "Импортировать данные в Cloud Firestore", fs=8)
uc["json"]    = oval(LX, r[8], LW, OH, "Подготовить JSON из Obsidian Vault", fs=8)

# ---- RIGHT COLUMN: include sub-cases ----
uc["search"]  = oval(RX, r[0], RW, OH, "Искать / фильтровать", fs=10)
uc["card"]    = oval(RX, r[1], RW, OH, "Открывать карточки",   fs=10)
uc["cr"]      = oval(RX, r[2], RW, OH, "Создать персонажа",    fs=10)
uc["ed"]      = oval(RX, r[3], RW, OH, "Редактировать персонажа", fs=9)
uc["sv"]      = oval(RX, r[4], RW, OH, "Сохранить персонажа",  fs=10)
uc["ld"]      = oval(RX, r[5], RW, OH, "Загрузить персонажа",  fs=10)
uc["del"]     = oval(RX, r[6], RW, OH, "Удалить персонажа",    fs=10)
uc["exp"]     = oval(RX, r[7], RW, OH, "Экспортировать персонажа", fs=9)
uc["impc"]    = oval(RX, r[8], RW, OH, "Импортировать персонажа",  fs=9)

# ---- INTERNAL ACTORS (left of boundary) ----
gost  = actor(0.47, (r[0] + r[4]) / 2,  "Гость",                      wl=0.90)
auth  = actor(0.47, (r[5] + r[6]) / 2,  "Авторизованный пользователь", wl=1.10)
admin = actor(0.47, (r[7] + r[8]) / 2,  "Администратор контента",      wl=1.10)

# ---- EXTERNAL SYSTEMS (right of boundary, cx < 10") ----
gauth = actor(9.10, 2.55, "Google Authentication", wl=1.30)
fs    = actor(9.10, 4.40, "Cloud Firestore",        wl=1.30)
ov    = actor(9.10, 6.10, "Obsidian Vault",         wl=1.30)

# ================================================================
# CONNECTORS
# ================================================================

def lft(k): return (uc[k]["cx"] - uc[k]["w"] / 2, uc[k]["cy"])
def rgt(k): return (uc[k]["cx"] + uc[k]["w"] / 2, uc[k]["cy"])


def assoc(a, k):
    """Solid association line: actor → use case (no arrowhead in UML)."""
    x1, y1 = a["cx"] + 0.28, a["cy"]
    x2, y2 = lft(k)
    connector(x1, y1, x2 - 0.02, y2, arrow=None)


# Гость → 5 use cases
for k in ("enter", "books", "news", "db", "login"):
    assoc(gost, k)

# Авторизованный пользователь → 2 use cases
for k in ("profile", "editor"):
    assoc(auth, k)

# Администратор → 2 use cases
for k in ("import", "json"):
    assoc(admin, k)


def include(src, dst, lbl=True):
    """Dashed include arrow."""
    x1, y1 = rgt(src)
    x2, y2 = lft(dst)
    connector(x1 + 0.02, y1, x2 - 0.02, y2, dash=True, arrow="stealth")
    if lbl:
        label((x1 + x2) / 2, (y1 + y2) / 2 - 0.13, "«include»")


include("db", "search")
include("db", "card")

# editor → 7 character-editor operations (one shared «include» label)
for k in ("cr", "ed", "sv", "ld", "del", "exp", "impc"):
    include("editor", k, lbl=False)
label((LX + RX) / 2, r[6] - 0.11, "«include»")


def ext(k, target):
    """Dashed line from use-case right edge → external system actor."""
    x1, y1 = rgt(k)
    x2, y2 = target["cx"] - 0.28, target["cy"]
    connector(x1 + 0.02, y1, x2, y2, dash=True, arrow="stealth")


ext("login",  gauth)   # login → Google Authentication
ext("import", fs)      # import → Cloud Firestore
ext("json",   ov)      # Obsidian JSON → Obsidian Vault

# ================================================================
prs.save(PPTX)
print("USE CASE BUILT ->", PPTX)
